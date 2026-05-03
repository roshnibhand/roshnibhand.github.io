from __future__ import annotations

import io
import re
from typing import Iterable, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from company_research_engine import NewsItem, ResearchReport, ResearchSection


INK = RGBColor(23, 32, 38)
MUTED = RGBColor(93, 108, 118)
ACCENT = RGBColor(15, 118, 110)
RUST = RGBColor(164, 63, 43)
BG = RGBColor(247, 248, 243)
WHITE = RGBColor(255, 255, 255)


def slugify_filename(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9]+", "_", value.strip()).strip("_").lower()
    return cleaned or "company_research"


def _set_background(slide, color=BG) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _add_textbox(slide, x, y, w, h, text: str, size: int = 20, bold: bool = False, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def _add_footer(slide, page_label: str) -> None:
    _add_textbox(slide, 0.55, 7.05, 7.6, 0.25, "AI-generated information. Verify important facts before relying on this deck.", 8, False, MUTED)
    box = _add_textbox(slide, 11.65, 7.05, 1.1, 0.25, page_label, 8, False, MUTED)
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _split_sentences(text: str, max_items: int = 5) -> List[str]:
    cleaned = re.sub(r"\s+", " ", text or "").strip()
    if not cleaned:
        return ["No reliable public information found."]
    parts = re.split(r"(?<=[.!?])\s+", cleaned)
    bullets = [part.strip(" -*") for part in parts if part.strip()]
    if len(bullets) <= 1:
        bullets = [cleaned]
    return bullets[:max_items]


def _add_bullets(slide, x, y, w, h, bullets: Iterable[str], size: int = 17) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(7)


def _add_chip(slide, x, y, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.15), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = ACCENT
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE


def build_company_research_pptx(report: ResearchReport) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank)
    _set_background(title_slide)
    _add_chip(title_slide, 0.65, 0.55, "AI-GENERATED RESEARCH")
    _add_textbox(title_slide, 0.65, 1.55, 8.4, 1.4, report.company.title, 42, True, INK)
    _add_textbox(
        title_slide,
        0.7,
        3.05,
        9.6,
        1.15,
        "Company Intelligence Snapshot for research, applications, networking, and business context.",
        22,
        False,
        MUTED,
    )
    _add_textbox(
        title_slide,
        0.7,
        5.65,
        10.8,
        0.55,
        "Generated with Gemini and Google Search grounding. Treat as a research starting point, not a verified source of truth.",
        14,
        False,
        RUST,
    )
    _add_footer(title_slide, "1")

    for index, section in enumerate(report.sections, start=2):
        slide = prs.slides.add_slide(blank)
        _set_background(slide)
        _add_chip(slide, 0.6, 0.45, section.short_label.upper())
        _add_textbox(slide, 0.6, 0.95, 11.6, 0.8, section.question, 28, True, INK)
        _add_bullets(slide, 0.85, 2.05, 11.25, 4.55, _split_sentences(section.answer), 18)
        _add_footer(slide, str(index))

    sources_slide = prs.slides.add_slide(blank)
    _set_background(sources_slide)
    _add_chip(sources_slide, 0.6, 0.45, "SOURCES")
    _add_textbox(sources_slide, 0.6, 0.95, 11.7, 0.65, "Grounded sources and verification notes", 28, True, INK)
    source_bullets = [
        f"{source.title}: {source.link}" for source in report.news[:7]
    ] or ["Gemini did not return source links for this search. Verify through company website, reputable news, investor pages, and LinkedIn."]
    _add_bullets(sources_slide, 0.85, 2.0, 11.5, 4.6, source_bullets, 14)
    _add_footer(sources_slide, str(len(report.sections) + 2))

    handle = io.BytesIO()
    prs.save(handle)
    return handle.getvalue()
